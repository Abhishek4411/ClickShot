# ClickShot - Button-only, fast screenshots with instant naming + PPTX/DOCX
# Windows 10/11 only. Run as Administrator recommended (for window picking precision).

import os, sys, time, threading, ctypes, queue
from ctypes import wintypes
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

import mss
from PIL import Image, ImageTk
import tkinter as tk
from tkinter import filedialog, simpledialog, messagebox

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from docx import Document
from docx.shared import Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH

import win32api, win32gui, win32con

APP_NAME = "ClickShot"

# ----- DPI awareness -----
user32 = ctypes.windll.user32
try:
    user32.SetProcessDpiAwarenessContext(ctypes.c_void_p(-4))  # Per Monitor v2
except Exception:
    pass

# DWM extended frame (accurate window bounds)
dwmapi = ctypes.windll.dwmapi
DWMWA_EXTENDED_FRAME_BOUNDS = 9

class RECT(ctypes.Structure):
    _fields_ = [('left', wintypes.LONG), ('top', wintypes.LONG),
                ('right', wintypes.LONG), ('bottom', wintypes.LONG)]

def get_extended_frame_bounds(hwnd):
    rect = RECT()
    try:
        dwmapi.DwmGetWindowAttribute(hwnd, DWMWA_EXTENDED_FRAME_BOUNDS,
                                     ctypes.byref(rect), ctypes.sizeof(rect))
        return (rect.left, rect.top, rect.right, rect.bottom)
    except Exception:
        return win32gui.GetWindowRect(hwnd)

# ----- Session -----
@dataclass
class Session:
    base_folder: Path
    project_name: str
    template_path: Path | None
    project_dir: Path
    pptx_path: Path
    docx_path: Path

# ----- Document Builder (background-safe) -----
class DocBuilder:
    def __init__(self, session: Session):
        self.session = session
        self.q = queue.Queue()
        self.t = threading.Thread(target=self._worker, daemon=True)
        self._running = True
        self.t.start()

    def init_docs(self): self.q.put(("init",))
    def add(self, image_path: Path, caption: str): self.q.put(("add", image_path, caption))
    def save(self): self.q.put(("save",))
    def close(self):
        self._running = False
        self.q.put(("close",))

    def _worker(self):
        prs, doc = None, None
        while self._running:
            try:
                msg = self.q.get(timeout=1)
            except queue.Empty:
                continue

            if msg[0] == "init":
                # PowerPoint
                try:
                    if self.session.template_path and self.session.template_path.exists():
                        prs = Presentation(str(self.session.template_path))
                    else:
                        prs = Presentation()
                    if len(prs.slides) == 0:
                        layout = prs.slide_layouts[0] if prs.slide_layouts else None
                        if layout:
                            slide = prs.slides.add_slide(layout)
                            if slide.shapes.title:
                                slide.shapes.title.text = self.session.project_name
                            if len(slide.placeholders) > 1:
                                try:
                                    slide.placeholders[1].text = f"Generated {datetime.now():%Y-%m-%d}"
                                except Exception:
                                    pass
                except Exception:
                    prs = Presentation()

                # Word
                try:
                    doc = Document()
                    h = doc.add_heading(self.session.project_name, 0)
                    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    p = doc.add_paragraph(f"Generated on {datetime.now():%Y-%m-%d %H:%M:%S}")
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    doc.add_page_break()
                except Exception:
                    doc = Document()

            elif msg[0] == "add" and prs and doc:
                _, image_path, caption = msg
                if not image_path.exists():
                    continue

                # PPT slide
                try:
                    slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1])
                    slide_w, slide_h = prs.slide_width, prs.slide_height
                    with Image.open(image_path) as im:
                        w, h = im.size
                        ar = w / h
                    max_w = slide_w * 0.88
                    max_h = slide_h * 0.74
                    if max_w / ar <= max_h:
                        pic_w, pic_h = max_w, max_w / ar
                    else:
                        pic_h, pic_w = max_h, max_h * ar
                    left = (slide_w - pic_w) / 2
                    top  = (slide_h - pic_h) / 2 - Inches(0.6)
                    slide.shapes.add_picture(str(image_path), left, top, width=pic_w, height=pic_h)

                    tb = slide.shapes.add_textbox(left, top + pic_h + Inches(0.25), pic_w, Inches(0.8))
                    tf = tb.text_frame; tf.clear()
                    p = tf.paragraphs[0]; p.text = caption
                    p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(20); p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 110, 210)
                except Exception:
                    pass

                # Word page
                try:
                    h1 = doc.add_heading(caption, level=1)
                    h1.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    par = doc.add_paragraph(); run = par.add_run()
                    try:
                        run.add_picture(str(image_path), width=DocxInches(6.5))
                    except Exception:
                        run.add_picture(str(image_path), width=DocxInches(6.0))
                    par.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    doc.add_page_break()
                except Exception:
                    pass

            elif msg[0] == "save":
                try:
                    if prs: prs.save(str(self.session.pptx_path))
                    if doc: doc.save(str(self.session.docx_path))
                except Exception:
                    pass

            elif msg[0] == "close":
                try:
                    if prs: prs.save(str(self.session.pptx_path))
                    if doc: doc.save(str(self.session.docx_path))
                except Exception:
                    pass
                break

# ----- Capture helpers -----
class Capture:
    def __init__(self):
        self.sct = mss.mss()

    @staticmethod
    def virtual_screen():
        SM_XVIRTUALSCREEN, SM_YVIRTUALSCREEN = 76, 77
        SM_CXVIRTUALSCREEN, SM_CYVIRTUALSCREEN = 78, 79
        x = user32.GetSystemMetrics(SM_XVIRTUALSCREEN)
        y = user32.GetSystemMetrics(SM_YVIRTUALSCREEN)
        w = user32.GetSystemMetrics(SM_CXVIRTUALSCREEN)
        h = user32.GetSystemMetrics(SM_CYVIRTUALSCREEN)
        return x, y, w, h

    @staticmethod
    def monitor_workarea_under_cursor():
        x, y = win32api.GetCursorPos()
        hmon = win32api.MonitorFromPoint((x, y), win32con.MONITOR_DEFAULTTONEAREST)
        info = win32api.GetMonitorInfo(hmon)
        l, t, r, b = info["Work"]  # excludes taskbar on that monitor
        return {"left": l, "top": t, "width": r - l, "height": b - t}

    def grab(self, box: dict) -> Image.Image:
        shot = self.sct.grab(box)
        return Image.frombytes("RGB", (shot.width, shot.height), shot.bgra, "raw", "BGRX")

    def capture_monitor(self) -> Image.Image:
        return self.grab(self.monitor_workarea_under_cursor())

# ----- Overlays -----
class RegionSelector:
    def __init__(self, root):
        self.root = root

    def select(self):
        """Drag rectangle over full virtual desktop. Returns (x, y, w, h) or None."""
        x0, y0, wv, hv = Capture.virtual_screen()
        overlay = tk.Toplevel(self.root)
        overlay.overrideredirect(True)
        overlay.attributes("-topmost", True)
        try: overlay.attributes("-alpha", 0.25)
        except Exception: pass
        overlay.configure(bg="#000")
        overlay.geometry(f"{wv}x{hv}+{x0}+{y0}")
        canvas = tk.Canvas(overlay, bg="#000", highlightthickness=0, cursor="crosshair")
        canvas.pack(fill="both", expand=True)

        canvas.create_text(wv//2, 40, text="Drag to select area • ESC to cancel",
                           fill="#58a6ff", font=("Segoe UI", 14, "bold"))

        result = {"box": None}
        rect = None
        start = {"x": 0, "y": 0}

        def on_press(e):
            nonlocal rect
            start["x"], start["y"] = e.x, e.y
            if rect: canvas.delete(rect)
            rect = canvas.create_rectangle(e.x, e.y, e.x, e.y, outline="#58a6ff", width=2)

        def on_drag(e):
            if rect: canvas.coords(rect, start["x"], start["y"], e.x, e.y)

        def on_release(e):
            if rect:
                x1, y1 = start["x"] + x0, start["y"] + y0
                x2, y2 = e.x + x0, e.y + y0
                if x2 < x1: x1, x2 = x2, x1
                if y2 < y1: y1, y2 = y2, y1
                if (x2 - x1) > 8 and (y2 - y1) > 8:
                    result["box"] = (x1, y1, x2 - x1, y2 - y1)
            overlay.destroy()

        def on_escape(_): overlay.destroy()

        canvas.bind("<Button-1>", on_press)
        canvas.bind("<B1-Motion>", on_drag)
        canvas.bind("<ButtonRelease-1>", on_release)
        overlay.bind("<Escape>", on_escape)

        overlay.focus_force()
        overlay.grab_set()
        self.root.update_idletasks()
        overlay.wait_window()  # no nested mainloop
        return result["box"]

class WindowPicker:
    def __init__(self, root):
        self.root = root

    def pick(self):
        """Click a window to capture. Returns hwnd or None."""
        x0, y0, wv, hv = Capture.virtual_screen()
        overlay = tk.Toplevel(self.root)
        overlay.overrideredirect(True)
        overlay.attributes("-topmost", True)
        try: overlay.attributes("-alpha", 0.15)
        except Exception: pass
        overlay.configure(bg="#000")
        overlay.geometry(f"{wv}x{hv}+{x0}+{y0}")
        canvas = tk.Canvas(overlay, bg="#000", highlightthickness=0, cursor="hand2")
        canvas.pack(fill="both", expand=True)
        canvas.create_text(wv//2, 40, text="Click a window to capture • ESC to cancel",
                           fill="#58a6ff", font=("Segoe UI", 14, "bold"))

        result = {"hwnd": None}

        def on_click(e):
            sx, sy = e.x + x0, e.y + y0
            hwnd = win32gui.WindowFromPoint((sx, sy))
            if hwnd:
                hwnd = win32gui.GetAncestor(hwnd, win32con.GA_ROOT)
                # Filter out desktop/taskbar/our own windows
                cname = win32gui.GetClassName(hwnd)
                title = win32gui.GetWindowText(hwnd) or ""
                if cname.lower() in ("shell_traywnd", "progman", "workerw"):
                    hwnd = None
                if title.find(APP_NAME) != -1:
                    hwnd = None
            result["hwnd"] = hwnd
            overlay.destroy()

        def on_escape(_):
            overlay.destroy()

        canvas.bind("<Button-1>", on_click)
        overlay.bind("<Escape>", on_escape)
        overlay.focus_force()
        overlay.grab_set()
        self.root.update_idletasks()
        overlay.wait_window()
        return result["hwnd"]

# ----- App -----
class ClickShotApp:
    def __init__(self):
        self.root = tk.Tk()
        self.root.title(APP_NAME)
        self.root.geometry("520x520")
        self.root.configure(bg="#0d1117")
        self.root.attributes("-alpha", 0.98)

        self.capture = Capture()
        self.selector = RegionSelector(self.root)
        self.winpicker = WindowPicker(self.root)

        self.session = self._setup_session()  # ask every run
        self.builder = DocBuilder(self.session)
        self.builder.init_docs()

        self._build_ui()
        self.root.protocol("WM_DELETE_WINDOW", self._quit)

    # -- Startup wizard --
    def _setup_session(self) -> Session:
        folder = filedialog.askdirectory(title="Choose a folder to save screenshots & documents")
        if not folder:
            messagebox.showerror(APP_NAME, "You must choose a save folder.")
            sys.exit(1)

        project = simpledialog.askstring(APP_NAME, "Project name:", initialvalue="ClickShot Project")
        if not project:
            project = "ClickShot_Project"

        template = None
        if messagebox.askyesno(APP_NAME, "Use a PowerPoint template (.pptx)?"):
            t = filedialog.askopenfilename(title="Choose PowerPoint Template", filetypes=[("PowerPoint", "*.pptx")])
            if t: template = Path(t)

        base = Path(folder)
        proj_dir = base / project.replace(" ", "_")
        proj_dir.mkdir(parents=True, exist_ok=True)
        return Session(
            base_folder=base,
            project_name=project,
            template_path=template,
            project_dir=proj_dir,
            pptx_path=proj_dir / f"{proj_dir.name}.pptx",
            docx_path=proj_dir / f"{proj_dir.name}.docx",
        )

    # -- UI --
    def _build_ui(self):
        fg, bg, glass, accent = "#f0f6fc", "#0d1117", "#161b22", "#58a6ff"

        wrap = tk.Frame(self.root, bg=bg, padx=18, pady=18); wrap.pack(fill="both", expand=True)

        header = tk.Frame(wrap, bg=glass); header.pack(fill="x", pady=(0, 16))
        tk.Label(header, text=APP_NAME, bg=glass, fg=accent, font=("Segoe UI", 22, "bold")).pack(pady=14)
        tk.Label(header, text="Buttons only • App hides during capture • Name dialog auto-focused",
                 bg=glass, fg=fg, font=("Segoe UI", 9, "italic")).pack(pady=(0, 10))

        status = tk.Frame(wrap, bg=glass); status.pack(fill="x", pady=(0, 16))
        self.status_lbl = tk.Label(status, text="Ready", bg=glass, fg="#16a34a",
                                   font=("Segoe UI", 11, "bold"), anchor="w", padx=12, pady=10)
        self.status_lbl.pack(fill="x")

        controls = tk.Frame(wrap, bg=glass); controls.pack(fill="x", pady=(0, 16))
        self._btn(controls, "📷 Capture Current Monitor (taskbar excluded)", self._capture_monitor).pack(pady=6, padx=12)
        self._btn(controls, "🪟 Capture Active Window (click target)", self._capture_window).pack(pady=6, padx=12)
        self._btn(controls, "✂️ Capture Region (drag rectangle)", self._capture_region).pack(pady=6, padx=12)

        proj = tk.Frame(wrap, bg=glass); proj.pack(fill="x", pady=(0, 16))
        tk.Label(proj, text="Project folder", bg=glass, fg=fg, font=("Segoe UI", 10, "bold"),
                 anchor="w").pack(fill="x", padx=12, pady=(10, 2))
        tk.Label(proj, text=str(self.session.project_dir), bg=glass, fg=fg, font=("Segoe UI", 9),
                 anchor="w", wraplength=460, justify="left").pack(fill="x", padx=12, pady=(0, 10))
        row = tk.Frame(proj, bg=glass); row.pack(pady=(0, 10))
        self._btn(row, "💾 Save PPTX/DOCX Now", lambda: self.builder.save()).pack(side="left", padx=6)
        self._btn(row, "📂 Change Project", self._change_project).pack(side="left", padx=6)

        footer = tk.Frame(wrap, bg=glass); footer.pack(fill="x")
        self._btn(footer, "❌ Quit", self._quit).pack(side="right", padx=12, pady=10)

    def _btn(self, parent, text, cmd):
        return tk.Button(parent, text=text, command=cmd, bd=0,
                         bg="#58a6ff", fg="white", activebackground="#1f6feb",
                         activeforeground="white", font=("Segoe UI", 10, "bold"),
                         padx=18, pady=12, cursor="hand2")

    # -- Status & Toasts --
    def _status(self, text, color): 
        try: self.status_lbl.configure(text=text, fg=color)
        except Exception: pass
    def _ok(self, text): self._status(text, "#16a34a")
    def _info(self, text): self._status(text, "#58a6ff")
    def _err(self, text): self._status(text, "#dc2626")

    def _toast(self, text):
        try:
            toast = tk.Toplevel(self.root)
            toast.overrideredirect(True)
            toast.attributes("-topmost", True)
            toast.attributes("-alpha", 0.0)
            toast.configure(bg="#161b22")
            tk.Label(toast, text=text, bg="#161b22", fg="#f0f6fc",
                     font=("Segoe UI", 9, "bold"), padx=18, pady=12).pack()
            toast.geometry("+110+110"); toast.update()
            for i in range(8): toast.attributes("-alpha", min(0.96, i*0.12)); toast.update(); time.sleep(0.015)
            toast.after(1600, lambda: self._fade_toast(toast))
        except Exception: pass

    def _fade_toast(self, toast):
        try:
            for i in range(8): toast.attributes("-alpha", max(0, 0.96 - i*0.12)); toast.update(); time.sleep(0.015)
            toast.destroy()
        except Exception: pass

    # -- App visibility helpers --
    def _hide_app(self):
        try:
            self.root.withdraw()
            self.root.update_idletasks()
            time.sleep(0.12)
        except Exception: pass

    def _show_app(self):
        try:
            self.root.deiconify()
            self.root.lift()
            self.root.attributes("-topmost", True)
            self.root.update_idletasks()
            self.root.after(120, lambda: self.root.attributes("-topmost", False))
        except Exception: pass

    # -- Capture handlers (buttons) --
    def _capture_monitor(self):
        self._hide_app()
        try:
            img = self.capture.capture_monitor()
            self._name_and_save(img, "monitor")
        except Exception as e:
            self._err(f"Capture failed: {e}")
        finally:
            self._show_app()

    def _capture_window(self):
        self._hide_app()
        try:
            hwnd = self.winpicker.pick()
            if not hwnd:
                self._info("Window capture cancelled.")
            else:
                l, t, r, b = get_extended_frame_bounds(hwnd)
                w, h = max(1, r - l), max(1, b - t)
                box = {"left": l, "top": t, "width": w, "height": h}
                img = self.capture.grab(box)
                self._name_and_save(img, "window")
        except Exception as e:
            self._err(f"Capture failed: {e}")
        finally:
            self._show_app()

    def _capture_region(self):
        self._hide_app()
        try:
            box = self.selector.select()
            if not box:
                self._info("Selection cancelled.")
            else:
                x, y, w, h = box
                img = self.capture.grab({"left": x, "top": y, "width": w, "height": h})
                self._name_and_save(img, "region")
        except Exception as e:
            self._err(f"Capture failed: {e}")
        finally:
            self._show_app()

    # -- Naming + saving --
    def _name_and_save(self, image: Image.Image, mode: str):
        name = self._name_dialog(image)  # dialog appears while app is hidden; entry is auto-focused
        if name is None:
            self._info("Save cancelled.")
            return

        safe = "".join(c for c in name if c.isalnum() or c in " -_()").strip()
        if not safe:
            safe = f"screenshot_{datetime.now():%H%M%S}"
        path = (self.session.project_dir / safe).with_suffix(".png")
        i = 1
        while path.exists():
            path = (self.session.project_dir / f"{safe}_{i}").with_suffix(".png")
            i += 1

        image.save(str(path), "PNG", optimize=True)
        self.builder.add(path, f"{safe} ({mode})")
        self.builder.save()

        # Bring app back on top and show toast
        self._show_app()
        self._ok(f"Saved: {path.name}")
        self._toast("✅ Saved successfully")

    def _name_dialog(self, pil_img: Image.Image) -> str | None:
        # Prepare preview
        preview = pil_img.copy()
        max_w, max_h = 540, 320
        w, h = preview.size
        ar = w / h
        if w > max_w or h > max_h:
            if max_w / ar <= max_h:
                nw, nh = int(max_w), int(max_w / ar)
            else:
                nh, nw = int(max_h), int(max_h * ar)
            preview = preview.resize((nw, nh), Image.LANCZOS)

        dlg = tk.Toplevel(self.root)
        dlg.title("Name your screenshot")
        dlg.configure(bg="#161b22")
        dlg.resizable(False, False)
        dlg.attributes("-topmost", True)
        dlg.grab_set()  # modal/focused
        dlg.focus_force()

        img_tk = ImageTk.PhotoImage(preview)
        tk.Label(dlg, image=img_tk, bg="#161b22").pack(padx=14, pady=(14, 8))

        tk.Label(dlg, text="File name (without extension):", bg="#161b22",
                 fg="#f0f6fc", font=("Segoe UI", 10, "bold")).pack(padx=14, anchor="w")

        entry = tk.Entry(dlg, width=50, font=("Segoe UI", 10),
                         bg="#0d1117", fg="#f0f6fc", insertbackground="#f0f6fc",
                         relief="flat")
        default_name = f"screenshot_{datetime.now():%Y%m%d_%H%M%S}"
        entry.insert(0, default_name)
        entry.pack(padx=14, pady=(6, 14))
        entry.select_range(0, tk.END)
        entry.focus_set()
        dlg.after(50, entry.focus_force)  # ensure focus sticks

        result = {"name": None}
        btns = tk.Frame(dlg, bg="#161b22"); btns.pack(pady=(0, 14))

        def ok():
            result["name"] = entry.get().strip()
            dlg.destroy()
        def cancel():
            dlg.destroy()

        tk.Button(btns, text="Save", command=ok, bg="#58a6ff", fg="white",
                  bd=0, padx=18, pady=8, font=("Segoe UI", 9, "bold"),
                  cursor="hand2").pack(side="left", padx=6)
        tk.Button(btns, text="Cancel", command=cancel, bg="#cf222e", fg="white",
                  bd=0, padx=18, pady=8, font=("Segoe UI", 9, "bold"),
                  cursor="hand2").pack(side="left", padx=6)

        dlg.bind("<Return>", lambda e: ok())
        dlg.bind("<Escape>", lambda e: cancel())

        self.root.update_idletasks()
        dlg.wait_window()  # no nested mainloop
        return result["name"]

    # -- Project switching & quit --
    def _change_project(self):
        self.session = self._setup_session()
        self.builder.close()
        self.builder = DocBuilder(self.session)
        self.builder.init_docs()
        self._ok("Project changed.")

    def _quit(self):
        try:
            self.builder.close()
        except Exception:
            pass
        finally:
            self.root.destroy()

    def run(self):
        self.root.mainloop()

# ----- Entrypoint -----
def ensure_deps():
    missing = []
    for mod, pipname in [("mss", "mss"), ("PIL", "Pillow"),
                         ("pptx", "python-pptx"), ("docx", "python-docx"),
                         ("win32api", "pywin32")]:
        try: __import__(mod)
        except Exception: missing.append(pipname)
    if missing:
        print("Missing modules:", ", ".join(missing))
        print("Install with:", "pip install " + " ".join(missing))
        input("Press Enter to exit...")
        sys.exit(1)

def main():
    if os.name != "nt":
        print(f"{APP_NAME} requires Windows.")
        return
    ensure_deps()
    ClickShotApp().run()

if __name__ == "__main__":
    main()
